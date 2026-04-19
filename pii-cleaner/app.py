import os
import uuid
import threading
import json
import zipfile
import io
from pathlib import Path
import tempfile

from flask import Flask, render_template, request, jsonify, send_file
from dotenv import load_dotenv
from openai import AzureOpenAI

# Text extraction
import pdfplumber
from docx import Document
from pptx import Presentation
import openpyxl

# PDF generation
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# Load .env from parent directory (Business Website Creation/.env)
load_dotenv(dotenv_path=Path(__file__).parent.parent / '.env')

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024  # 200 MB max upload

AZURE_KEY        = os.getenv('AZURE_OPENAI_API_KEY')
AZURE_ENDPOINT   = os.getenv('AZURE_OPENAI_ENDPOINT')
AZURE_DEPLOYMENT = os.getenv('AZURE_OPENAI_DEPLOYMENT', 'gpt-4o-mini')
AZURE_API_VER    = os.getenv('AZURE_OPENAI_API_VERSION', '2024-12-01-preview')

TEMP_DIR = Path(tempfile.gettempdir()) / 'pii_cleaner'
TEMP_DIR.mkdir(exist_ok=True)

jobs: dict = {}
jobs_lock = threading.Lock()

# ─── Register a Unicode-capable font for PDF output ─────────────────────────
BODY_FONT = 'Helvetica'
_font_candidates = [
    '/Library/Fonts/Arial Unicode MS.ttf',
    '/Library/Fonts/Arial.ttf',
    '/System/Library/Fonts/Supplemental/Arial.ttf',
    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    'C:/Windows/Fonts/arial.ttf',
]
for _fp in _font_candidates:
    if Path(_fp).exists():
        try:
            pdfmetrics.registerFont(TTFont('UniFont', _fp))
            BODY_FONT = 'UniFont'
        except Exception:
            pass
        break

# ─── System prompt ───────────────────────────────────────────────────────────
SYSTEM_PROMPT = """You are a data anonymization expert. Your task is to find and replace ALL personal data in the provided text with appropriate placeholders. Be thorough and miss nothing.

Replace the following types of personal data:
- Full names or partial names of real people → [NAME]
- Specific dates (birthdays, appointment dates, etc.) → [DATE]
- Phone numbers (in any format) → [PHONE]
- Email addresses → [EMAIL]
- Physical addresses (street, house number, city, postal code, country) → [ADDRESS]
- Company names and business names → [COMPANY]
- School names, university names, educational institution names → [SCHOOL/UNIVERSITY]

Rules:
- Replace ONLY the personal data, keep all other text exactly as-is
- Be thorough — catch every single instance
- Preserve the exact formatting, line breaks, and spacing of the original text
- Return ONLY the anonymized text with absolutely no commentary, explanation, or preamble"""


# ─── Text extraction ─────────────────────────────────────────────────────────
def extract_text(filepath: Path) -> str:
    ext = filepath.suffix.lower()

    if ext == '.pdf':
        pages = []
        with pdfplumber.open(str(filepath)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    pages.append(t)
        return '\n\n'.join(pages)

    elif ext == '.docx':
        doc = Document(str(filepath))
        return '\n'.join(p.text for p in doc.paragraphs)

    elif ext == '.pptx':
        prs = Presentation(str(filepath))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f'[Slide {i}]')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    parts.append(shape.text)
        return '\n'.join(parts)

    elif ext == '.xlsx':
        wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            rows.append(f'[Sheet: {ws.title}]')
            for row in ws.iter_rows(values_only=True):
                row_text = '\t'.join(str(c) if c is not None else '' for c in row)
                if row_text.strip():
                    rows.append(row_text)
        return '\n'.join(rows)

    elif ext == '.rtf':
        from striprtf.striprtf import rtf_to_text
        with open(str(filepath), 'r', encoding='utf-8', errors='ignore') as f:
            return rtf_to_text(f.read())

    else:
        # Treat as plain text (txt, csv, md, log, html, json, xml, etc.)
        with open(str(filepath), 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()


# ─── PII cleaning via Azure OpenAI ──────────────────────────────────────────
def clean_pii(text: str) -> str:
    client = AzureOpenAI(
        api_version=AZURE_API_VER,
        azure_endpoint=AZURE_ENDPOINT,
        api_key=AZURE_KEY,
    )

    def call_api(chunk: str) -> str:
        resp = client.chat.completions.create(
            model=AZURE_DEPLOYMENT,
            messages=[
                {'role': 'system', 'content': SYSTEM_PROMPT},
                {'role': 'user',   'content': chunk},
            ],
            max_tokens=16000,
            temperature=0.0,
        )
        return resp.choices[0].message.content or ''

    MAX_CHUNK = 12_000
    if len(text) <= MAX_CHUNK:
        return call_api(text)

    # Split on paragraph boundaries to avoid cutting mid-sentence
    lines = text.split('\n')
    chunks, current, current_len = [], [], 0
    for line in lines:
        line_len = len(line) + 1
        if current_len + line_len > MAX_CHUNK and current:
            chunks.append('\n'.join(current))
            current, current_len = [line], line_len
        else:
            current.append(line)
            current_len += line_len
    if current:
        chunks.append('\n'.join(current))

    return '\n'.join(call_api(c) for c in chunks)


# ─── PDF generation ──────────────────────────────────────────────────────────
def text_to_pdf(text: str, output_path: Path) -> None:
    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        rightMargin=22 * mm,
        leftMargin=22 * mm,
        topMargin=20 * mm,
        bottomMargin=20 * mm,
    )
    styles = getSampleStyleSheet()
    body_style = ParagraphStyle(
        'body',
        parent=styles['Normal'],
        fontName=BODY_FONT,
        fontSize=10,
        leading=15,
        spaceAfter=3,
        wordWrap='LTR',
    )

    story = []
    for line in text.split('\n'):
        safe = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        if safe.strip():
            story.append(Paragraph(safe, body_style))
        else:
            story.append(Spacer(1, 7))

    if not story:
        story.append(Paragraph('[No content after anonymization]', body_style))

    doc.build(story)


# ─── Background processing job ───────────────────────────────────────────────
def process_job(job_id: str, file_paths: list, original_names: list) -> None:
    out_dir = TEMP_DIR / job_id / 'output'
    out_dir.mkdir(parents=True, exist_ok=True)

    with jobs_lock:
        jobs[job_id].update(status='processing', total=len(file_paths), done=0)

    for i, (fpath, fname) in enumerate(zip(file_paths, original_names)):
        stem     = Path(fname).stem
        out_name = f'{stem}_cleaned.pdf'
        out_path = out_dir / out_name

        with jobs_lock:
            jobs[job_id]['current_file'] = fname

        try:
            text = extract_text(Path(fpath))
            if not text.strip():
                text = '[No extractable text found in this file]'

            cleaned = clean_pii(text)
            text_to_pdf(cleaned, out_path)

            with jobs_lock:
                jobs[job_id]['outputs'].append({
                    'name':   out_name,
                    'path':   str(out_path),
                    'source': fname,
                })
        except Exception as exc:
            with jobs_lock:
                jobs[job_id]['errors'].append({'file': fname, 'error': str(exc)})

        with jobs_lock:
            jobs[job_id]['done'] = i + 1

    with jobs_lock:
        jobs[job_id].update(status='done', current_file='')


# ─── Routes ──────────────────────────────────────────────────────────────────
@app.route('/')
def index():
    return render_template('index.html',
                           config_ok=bool(AZURE_KEY and AZURE_ENDPOINT))


@app.route('/process', methods=['POST'])
def process():
    if 'files' not in request.files:
        return jsonify({'error': 'No files provided'}), 400

    files = [f for f in request.files.getlist('files') if f.filename]
    if not files:
        return jsonify({'error': 'No valid files'}), 400

    job_id    = str(uuid.uuid4())
    input_dir = TEMP_DIR / job_id / 'input'
    input_dir.mkdir(parents=True, exist_ok=True)

    file_paths, original_names = [], []
    for f in files:
        dest = input_dir / f.filename
        f.save(str(dest))
        file_paths.append(str(dest))
        original_names.append(f.filename)

    with jobs_lock:
        jobs[job_id] = {
            'status':       'queued',
            'total':        len(file_paths),
            'done':         0,
            'current_file': '',
            'outputs':      [],
            'errors':       [],
        }

    threading.Thread(
        target=process_job,
        args=(job_id, file_paths, original_names),
        daemon=True,
    ).start()

    return jsonify({'job_id': job_id, 'file_count': len(file_paths)})


@app.route('/status/<job_id>')
def status(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return jsonify({'error': 'Job not found'}), 404
    return jsonify(job)


@app.route('/download/<job_id>/<path:filename>')
def download_file(job_id, filename):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return 'Job not found', 404
    for out in job['outputs']:
        if out['name'] == filename:
            return send_file(out['path'], as_attachment=True, download_name=filename)
    return 'File not found', 404


@app.route('/download-zip/<job_id>')
def download_zip(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job or not job['outputs']:
        return 'Nothing to download', 404

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for out in job['outputs']:
            zf.write(out['path'], out['name'])
    buf.seek(0)
    return send_file(
        buf,
        as_attachment=True,
        download_name='cleaned_documents.zip',
        mimetype='application/zip',
    )


# ─── Entry point ─────────────────────────────────────────────────────────────
if __name__ == '__main__':
    import subprocess
    import time

    def open_chrome():
        time.sleep(1.5)
        try:
            subprocess.Popen(['open', '-a', 'Google Chrome', 'http://127.0.0.1:5055'])
        except Exception:
            import webbrowser
            webbrowser.open('http://127.0.0.1:5055')

    threading.Thread(target=open_chrome, daemon=True).start()

    print('\n  PII Cleaner  →  http://127.0.0.1:5055\n')
    app.run(port=5055, debug=False, threaded=True)
